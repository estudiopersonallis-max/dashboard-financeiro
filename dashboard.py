# ================= IMPORTS =================
import streamlit as st
import pandas as pd
import matplotlib.pyplot as plt
from io import BytesIO
from datetime import datetime
import re
import unicodedata

# PDF
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Image, PageBreak
from reportlab.lib.pagesizes import A4
from reportlab.lib.styles import getSampleStyleSheet
from reportlab.lib.units import cm

# PPTX
from pptx import Presentation
from pptx.util import Inches
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE

# ================= CONFIG =================
st.set_page_config(page_title="Dashboard Financeiro PRO", layout="wide")
st.title("📊 Dashboard Financeiro – Nível Consultoria")

# ================= NORMALIZAÇÃO =================
def normalizar(txt):
    if pd.isna(txt):
        return ""
    txt = str(txt).upper().strip()
    txt = unicodedata.normalize('NFKD', txt).encode('ASCII', 'ignore').decode('ASCII')
    return txt

# ================= DETECTAR MÊS =================
mapa_meses = {
    "JAN":1, "JANEIRO":1,
    "FEV":2, "FEVEREIRO":2,
    "MAR":3, "MARCO":3,
    "ABR":4, "ABRIL":4,
    "MAI":5, "MAIO":5,
    "JUN":6, "JUNHO":6,
    "JUL":7, "JULHO":7,
    "AGO":8, "AGOSTO":8,
    "SET":9, "SETEMBRO":9,
    "OUT":10, "OUTUBRO":10,
    "NOV":11, "NOVEMBRO":11,
    "DEZ":12, "DEZEMBRO":12
}

def extrair_mes(nome):
    nome = normalizar(nome)

    match = re.search(r'\b(0?[1-9]|1[0-2])\b', nome)
    if match:
        return int(match.group())

    for k, v in mapa_meses.items():
        if k in nome:
            return v

    return 99

# ================= LEITURA =================
@st.cache_data(ttl=3600)
def ler_receitas(files):
    dfs = []
    for f in files:
        df = pd.read_excel(f)
        if df.empty:
            continue

        periodo_nome = f.name.split(".")[0]
        mes = extrair_mes(periodo_nome)

        df["Periodo"] = periodo_nome.upper()
        df["ordem_mes"] = mes

        df["Valor"] = pd.to_numeric(df.get("Valor", 0), errors="coerce").fillna(0)

        df["Nome do cliente"] = df.get("Nome do cliente", "").apply(normalizar)
        df = df[df["Nome do cliente"] != ""]

        df["Modalidade"] = df.get("Modalidade", "N/A").apply(normalizar)
        df["Modalidade"] = df["Modalidade"].replace("", "SEM MODALIDADE")

        df["Tipo"] = df.get("Tipo", "N/A")
        df["Professor"] = df.get("Professor", "N/A")
        df["Local"] = df.get("Local", "N/A")

        dfs.append(df)

    return pd.concat(dfs, ignore_index=True) if dfs else pd.DataFrame()

@st.cache_data(ttl=3600)
def ler_despesas(files):
    dfs = []
    for f in files:
        df = pd.read_excel(f)
        df = df.dropna(subset=["Valor", "Descrição da Despesa", "Classe"])
        if df.empty:
            continue

        periodo_nome = f.name.split(".")[0]
        mes = extrair_mes(periodo_nome)

        df["Periodo"] = periodo_nome.upper()
        df["ordem_mes"] = mes

        df["Valor"] = pd.to_numeric(df.get("Valor", 0), errors="coerce").fillna(0)
        df["Classe"] = df.get("Classe", "N/A").apply(normalizar)

        df["Local"] = df.get("Local", "N/A")

        dfs.append(df)

    return pd.concat(dfs, ignore_index=True) if dfs else pd.DataFrame()

# ================= FUNÇÕES AUX =================
def grafico_bar(df, titulo):
    fig, ax = plt.subplots()
    df.plot(kind="barh", ax=ax)
    ax.set_title(titulo)
    return fig

def gerar_pivot(df, index):
    return df.pivot_table(
        index=index,
        columns="Periodo",
        values="Valor",
        aggfunc="sum",
        fill_value=0
    )

# ================= EXPORT HELPERS =================
figs_pdf = []

def capturar_grafico(fig, titulo, pivot):
    figs_pdf.append((titulo, fig, pivot))

# ================= UPLOAD =================
st.sidebar.header("📤 Upload")
uploaded_receitas = st.sidebar.file_uploader("Receitas", type=["xlsx"], accept_multiple_files=True)
uploaded_despesas = st.sidebar.file_uploader("Despesas", type=["xlsx"], accept_multiple_files=True)

receitas = ler_receitas(uploaded_receitas) if uploaded_receitas else pd.DataFrame()
despesas = ler_despesas(uploaded_despesas) if uploaded_despesas else pd.DataFrame()

# ================= FILTROS =================
st.sidebar.header("🔎 Filtros")
periodos = sorted(set(receitas.get("Periodo", [])).union(set(despesas.get("Periodo", []))))
periodo_sel = st.sidebar.multiselect("Períodos", periodos, default=periodos)

if not receitas.empty:
    receitas = receitas[receitas["Periodo"].isin(periodo_sel)]

if not despesas.empty:
    despesas = despesas[despesas["Periodo"].isin(periodo_sel)]
    despesas = despesas[despesas["Classe"] != "DEPOSITOS"]

# ================= KPIs =================
receita_total = receitas["Valor"].sum() if not receitas.empty else 0
despesa_total = despesas["Valor"].sum() if not despesas.empty else 0
lucro_total = receita_total + despesa_total
margem = (lucro_total / receita_total * 100) if receita_total else 0

receita_media = receitas.groupby("Periodo")["Valor"].sum().mean() if not receitas.empty else 0
despesa_media = despesas.groupby("Periodo")["Valor"].sum().mean() if not despesas.empty else 0
clientes_ativos_media = receitas.groupby("Periodo")["Nome do cliente"].nunique().mean() if not receitas.empty else 0

ticket_medio_receita = receita_media / clientes_ativos_media if clientes_ativos_media else 0
ticket_medio_despesa = abs(despesa_media) / clientes_ativos_media if clientes_ativos_media else 0

magic_number = abs(despesa_media)

cac = abs(despesa_media) / clientes_ativos_media if clientes_ativos_media else 0
ltv = ticket_medio_receita * 60 if ticket_medio_receita else 0
ltv_cac = (ltv / cac) if cac else 0

# KPIs em colunas
col1, col2, col3, col4 = st.columns(4)
col1.metric("Receita", f"{receita_total:,.0f}€")
col2.metric("Despesa", f"{despesa_total:,.0f}€")
col3.metric("Lucro", f"{lucro_total:,.0f}€")
col4.metric("Margem", f"{margem:.1f}%")

col5, col6, col7, col8 = st.columns(4)
col5.metric("Ticket Receita", f"{ticket_medio_receita:,.0f}€")
col6.metric("Ticket Despesa", f"{ticket_medio_despesa:,.0f}€")
col7.metric("CAC", f"{cac:,.0f}€")
col8.metric("LTV/CAC", f"{ltv_cac:.2f}")

# ================= ALERTAS =================
st.subheader("⚠️ Alertas Inteligentes")
if margem < 20:
    st.warning("Margem abaixo de 20%")
if despesa_total < -receita_total * 0.8:
    st.warning("Despesas muito altas")
if clientes_ativos_media < 10:
    st.warning("Poucos clientes ativos")

# ================= EVOLUÇÃO FINANCEIRA =================
st.subheader("📈 Evolução Financeira")
if not receitas.empty and not despesas.empty:
    receita_mes = receitas.groupby("Periodo")["Valor"].sum()
    despesa_mes = despesas.groupby("Periodo")["Valor"].sum()

    df_fin = pd.DataFrame({
        "Receita": receita_mes,
        "Despesa": despesa_mes
    }).fillna(0)

    df_fin["Lucro"] = df_fin["Receita"] + df_fin["Despesa"]

    fig, ax = plt.subplots()
    df_fin.plot(ax=ax, marker="o")
    plt.xticks(rotation=45)
    st.pyplot(fig)

# ================= TOP CLIENTES =================
st.subheader("🏆 Top Clientes")
if not receitas.empty:
    top_clientes = (
        receitas.groupby("Nome do cliente")["Valor"]
        .sum()
        .sort_values(ascending=False)
        .head(10)
    )
    st.dataframe(top_clientes)

# ================= CLIENTES =================
st.subheader("👥 Evolução de Clientes")
if not receitas.empty:
    clientes_por_mes = (
        receitas.groupby(["Periodo", "ordem_mes"])["Nome do cliente"]
        .nunique()
        .reset_index()
        .sort_values("ordem_mes")
    )

    fig, ax = plt.subplots()
    ax.plot(clientes_por_mes["Periodo"], clientes_por_mes["Nome do cliente"], marker="o")
    plt.xticks(rotation=45)
    st.pyplot(fig)

# ================= MODALIDADE =================
st.subheader("🏋️ Clientes por Modalidade")
if not receitas.empty:
    clientes_modalidade = receitas.groupby("Modalidade")["Nome do cliente"].nunique().sort_values(ascending=False)
    st.dataframe(clientes_modalidade)

    fig_mod, ax_mod = plt.subplots()
    clientes_modalidade.plot(kind="barh", ax=ax_mod)
    st.pyplot(fig_mod)

# ================= TABS =================
tab1, tab2, tab3 = st.tabs(["📊 Visão Geral", "💰 Receitas", "💸 Despesas"])

with tab2:
    for cat in ["Modalidade", "Tipo", "Professor", "Local"]:
        bloco = gerar_pivot(receitas, cat)
        st.dataframe(bloco)

        fig = grafico_bar(bloco, f"Receitas por {cat}")
        st.pyplot(fig)

        capturar_grafico(fig, f"Receitas por {cat}", bloco)

with tab3:
    for cat in ["Classe", "Local"]:
        bloco = gerar_pivot(despesas, cat)
        st.dataframe(bloco)

        fig = grafico_bar(bloco, f"Despesas por {cat}")
        st.pyplot(fig)

        capturar_grafico(fig, f"Despesas por {cat}", bloco)

# ================= EXPORT =================
st.subheader("📄 Exportações")

def gerar_pdf(figs):
    buffer = BytesIO()
    doc = SimpleDocTemplate(buffer, pagesize=A4)
    styles = getSampleStyleSheet()
    elementos = []

    for titulo, fig, _ in figs:
        elementos.append(Paragraph(titulo, styles["Heading2"]))

        img = BytesIO()
        fig.savefig(img, format="png", bbox_inches="tight")
        img.seek(0)

        elementos.append(Image(img, width=16*cm, height=8*cm))
        elementos.append(PageBreak())

    doc.build(elementos)
    buffer.seek(0)
    return buffer


def gerar_ppt(figs):
    prs = Presentation()

    for titulo, _, pivot in figs:
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        slide.shapes.title.text = titulo

        chart_data = CategoryChartData()
        chart_data.categories = list(pivot.index)

        for col in pivot.columns:
            chart_data.add_series(str(col), list(pivot[col].values))

        slide.shapes.add_chart(
            XL_CHART_TYPE.BAR_CLUSTERED,
            Inches(1), Inches(2), Inches(8), Inches(4),
            chart_data
        )

    buffer = BytesIO()
    prs.save(buffer)
    buffer.seek(0)
    return buffer

col1, col2 = st.columns(2)

with col1:
    if st.button("📄 Gerar PDF Completo"):
        st.download_button("Download PDF", gerar_pdf(figs_pdf), "relatorio.pdf")

with col2:
    if st.button("📊 Gerar PPT Editável"):
        st.download_button("Download PPT", gerar_ppt(figs_pdf), "relatorio.pptx")

# ================= FOOTER =================
st.caption(f"Atualizado em {datetime.now().strftime('%d/%m/%Y %H:%M')}")
